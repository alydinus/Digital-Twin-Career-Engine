from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # ---------------------------------------------------------
    # Helper to add a formatted title
    # ---------------------------------------------------------
    def add_slide(layout_idx, title_text):
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            # Customize title font
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 51, 102)
        return slide

    def add_bullets(slide, content_list):
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()  # Clear existing default text
            
            for item in content_list:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(20)

    # ---------------------------------------------------------
    # Slide 1: Title Slide
    # ---------------------------------------------------------
    slide0 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide0.shapes.title
    subtitle = slide0.placeholders[1]

    title.text = "Digital Twin Career Engine"
    
    # Make title bold and colored
    for run in title.text_frame.paragraphs[0].runs:
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 102, 204)

    subtitle.text = "Final Generative AI Project\n\nАвтор: Alydin\n2026"

    # ---------------------------------------------------------
    # Slide 2: О проекте
    # ---------------------------------------------------------
    s2 = add_slide(1, "О проекте: Digital Twin Career Engine")
    add_bullets(s2, [
        "Цель: Создать интеллектуального карьерного ИИ-ассистента.",
        "Анализ профиля: Извлечение навыков из CV, LinkedIn и NotebookLM.",
        "Карьерный трек: Предсказание идеальных ролей с помощью ML.",
        "Персонализация: План развития (Roadmap) и вопросы для интервью.",
        "GenAI: Интеграция с LLM (Anthropic / OpenAI / Gemini).",
        "Строгая 5-уровневая архитектура по требованиям проекта."
    ])

    # ---------------------------------------------------------
    # Slide 3: 5-Слойная Архитектура
    # ---------------------------------------------------------
    s3 = add_slide(1, "Архитектура Системы (5 Слоёв)")
    add_bullets(s3, [
        "1. Platform Layer: Интеграция данных (PDF, LinkedIn, MCP).",
        "2. Model Layer: Классический ML (Косинусное сходство, TF-IDF).",
        "3. Agent Layer: Автономные ИИ-агенты и Web Browsing.",
        "4. Application Layer: Streamlit UI (8 вкладок, виджеты).",
        "5. Infrastructure Layer: Локальное/Облачное выполнение, API."
    ])

    # ---------------------------------------------------------
    # Slide 4: Platform Layer
    # ---------------------------------------------------------
    s4 = add_slide(1, "Слой 1: Platform Layer (Интеграция Данных)")
    add_bullets(s4, [
        "Универсальный парсер профиля:",
        "  • PDF Резюме (pdfplumber + regex/LLM)",
        "  • LinkedIn Import (Поддержка PDF экспорта и текста)",
        "Интеграция с Google NotebookLM:",
        "  • Использование протокола MCP (Model Context Protocol).",
        "  • Извлечение базы знаний напрямую через notebooklm-mcp-server.",
        "Датасет: Расширенная база 21 профессии × 8+ навыков."
    ])

    # ---------------------------------------------------------
    # Slide 5: Model Layer
    # ---------------------------------------------------------
    s5 = add_slide(1, "Слой 2: Model Layer (Алгоритмы)")
    add_bullets(s5, [
        "Классический Machine Learning (scikit-learn):",
        "  • TF-IDF векторизация навыков.",
        "  • Cosine Similarity (Косинусное сходство).",
        "Метрики совпадения:",
        "  • Hybrid Match Score: комбинация семантики и точного пересечения.",
        "Нормализация:",
        "  • Приведение синонимов (k8s -> kubernetes) для точности."
    ])

    # ---------------------------------------------------------
    # Slide 6: Agent Layer
    # ---------------------------------------------------------
    s6 = add_slide(1, "Слой 3: Agent Layer (GenAI & Web)")
    add_bullets(s6, [
        "Автономные агенты с доступом в интернет:",
        "  • Web Searcher: DuckDuckGo парсинг + GitHub Search API.",
        "  • Resource Finder: Динамический поиск курсов и документации.",
        "Генеративные агенты:",
        "  • Career Coach: Пошаговый план обучения (Roadmap).",
        "  • Interview Coach: Подготовка вопросов по 'слабым местам'.",
        "  • Live Coach: Чат-бот с переключением персон."
    ])

    # ---------------------------------------------------------
    # Slide 7: Application Layer
    # ---------------------------------------------------------
    s7 = add_slide(1, "Слой 4: Application Layer (UI)")
    add_bullets(s7, [
        "Streamlit UI с 8 вкладками.",
        "Уникальные виджеты:",
        "  • Balance Wheel: Радар-чарт (Hard vs Soft skills).",
        "  • RPG Tech Tree: Дерево навыков в стиле видеоигры (Locked/Unlocked).",
        "  • Semester Wrapped: Экспорт карточки-достижения для LinkedIn.",
        "  • Live Coach: Интерактивный чат с режимом 'Roast My Stack'.",
        "  • Telegram Jobs: Просмотр подходящих вакансий."
    ])

    # ---------------------------------------------------------
    # Slide 8: Демо-сценарий
    # ---------------------------------------------------------
    s8 = add_slide(1, "Демонстрация (User Flow)")
    add_bullets(s8, [
        "1. Ввод данных: Загрузка PDF / LinkedIn / NotebookLM MCP.",
        "2. Predict: ML рассчитывает процент совпадения с 21 профессией.",
        "3. Аналитика: Изучение RPG Tech Tree и Balance Wheel.",
        "4. Общение с ИИ: 'Прожарка' (Roast) стека в Live Coach.",
        "5. Планирование: Генерация Roadmap и тренировка Interview.",
        "6. Результат: Поиск работы (Telegram Jobs) и скачивание Semester Wrapped."
    ])

    # ---------------------------------------------------------
    # Slide 9: Заключение
    # ---------------------------------------------------------
    s9 = add_slide(1, "Заключение")
    add_bullets(s9, [
        "Полное соответствие требованиям Final Generative AI Project.",
        "Реализованы все 5 слоёв архитектуры.",
        "Использованы современные подходы:",
        "  • MCP (Model Context Protocol).",
        "  • Retrieval-Augmented Generation (RAG).",
        "  • Agentic AI (Автономный поиск).",
        "Проект готов к деплою на Streamlit Community Cloud."
    ])

    # Save presentation
    prs.save("Digital_Twin_Presentation.pptx")
    print("Presentation saved as Digital_Twin_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
